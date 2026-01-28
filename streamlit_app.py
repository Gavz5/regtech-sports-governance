import os
import re
import json
import io
import time
from datetime import datetime
from typing import List, Dict, Any, Tuple

import pandas as pd
import requests
import streamlit as st

# Optional file extractors
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


# =========================
# CONFIG
# =========================
APP_TITLE = "RegTech Analytics for Sports Governance Failure"
DATA_DIR = os.path.dirname(os.path.abspath(__file__))
ATHLETES_FILE = os.path.join(DATA_DIR, "athletes_structured.json")
POLICY_CACHE_DIR = os.path.join(DATA_DIR, "policy_cache")
os.makedirs(POLICY_CACHE_DIR, exist_ok=True)

st.set_page_config(page_title=APP_TITLE, layout="wide")


# =========================
# UI STYLE (Locked feel)
# =========================
st.markdown(
    """
    <style>
      .block-container {padding-top: 1.2rem; padding-bottom: 2rem;}
      .stTabs [data-baseweb="tab-list"] {gap: 8px;}
      .stTabs [data-baseweb="tab"] {padding-left: 18px; padding-right: 18px;}
      .small-note {font-size: 0.92rem; opacity: 0.85;}
      .metric-card {padding: 14px; border-radius: 12px; border: 1px solid rgba(0,0,0,0.08);}
      .muted {opacity: 0.75;}
      .danger {color: #b00020;}
      .good {color: #0b6b2f;}
      .warn {color: #8a5a00;}
      code {white-space: pre-wrap !important;}
    </style>
    """,
    unsafe_allow_html=True
)


# =========================
# TAXONOMY (Policy obligations)
# =========================
OBLIGATIONS = [
    {
        "obligation": "Identity & Registration",
        "why": "Without a verified athlete identity + registration workflow, benefits and selection become arbitrary, leading to exclusion.",
        "keywords": ["athlete id", "registration", "unique id", "identity", "database", "profile", "enrolment", "enrollment", "verification"],
    },
    {
        "obligation": "Eligibility Criteria",
        "why": "Eligibility rules decide who gets access. If vague, it creates bias, favoritism, and denial without reason.",
        "keywords": ["eligibility", "criteria", "qualify", "qualification", "age", "category", "selection", "trial", "benchmarks"],
    },
    {
        "obligation": "Benefits & Entitlements",
        "why": "If the policy doesn't clearly list entitlements (scholarship, hostel, kit, insurance), athletes get denied silently.",
        "keywords": ["benefit", "entitlement", "scholarship", "stipend", "hostel", "kit", "equipment", "travel", "diet", "insurance", "medical"],
    },
    {
        "obligation": "Process Definition",
        "why": "Processes define how to apply, approve, and deliver support. Missing processes = delays + corruption risk.",
        "keywords": ["process", "procedure", "apply", "application", "workflow", "submission", "documents", "committee", "approval"],
    },
    {
        "obligation": "Timelines & Service Levels",
        "why": "Without timelines, authorities delay indefinitely; athletes lose opportunities (tournaments/admissions).",
        "keywords": ["timeline", "within", "days", "deadline", "time limit", "service level", "sla", "period"],
    },
    {
        "obligation": "Accountability & Ownership",
        "why": "Who is responsible? Without named responsibility, complaints go nowhere and denial has no consequence.",
        "keywords": ["responsible", "accountable", "nodal officer", "authority", "department", "sports officer", "principal", "director", "roles"],
    },
    {
        "obligation": "Grievance, Appeal & Remedies",
        "why": "If grievance & appeal are weak, athletes cannot challenge unfair denial—governance failure becomes permanent.",
        "keywords": ["grievance", "appeal", "complaint", "redressal", "remedy", "hearing", "review", "ombudsman"],
    },
    {
        "obligation": "Digitization & Transparency",
        "why": "Digitization improves traceability: who applied, who approved, when, and why. Without it, records get 'lost'.",
        "keywords": ["online", "portal", "digital", "dashboard", "transparency", "public", "track", "status", "sms", "email", "upload"],
    },
    {
        "obligation": "Enforcement & Penalties",
        "why": "Enforcement clauses force compliance. Without penalties, policy remains a suggestion.",
        "keywords": ["penalty", "disciplinary", "punishment", "fine", "action", "non-compliance", "audit", "inspection"],
    },
]


# Athlete-side entitlement choices (used in risk simulation)
ENTITLEMENT_CHOICES = [
    "Hostel",
    "Scholarship",
    "Sports Kit",
    "Travel Allowance",
    "Medical/Insurance",
    "Diet/Nutrition Support",
    "Fee Waiver",
    "Coaching Support",
]


# =========================
# HELPERS: Safe read/write
# =========================
def now_iso() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def safe_load_json(path: str, default):
    if not os.path.exists(path):
        return default
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return default


def safe_save_json(path: str, data):
    tmp = path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    os.replace(tmp, path)


# =========================
# ATHLETE DATA: normalize schema
# =========================
def normalize_athlete_record(rec: Dict[str, Any]) -> Dict[str, Any]:
    """
    Enforces ONE clean schema so you never get messy scattered columns.
    Also auto-fixes older keys like "entitlements" into "entitlements_claimed".
    """
    fixed = {}

    # Required core fields
    fixed["athlete_name"] = str(rec.get("athlete_name") or rec.get("full_name") or rec.get("name") or "").strip()
    fixed["sport"] = str(rec.get("sport") or "").strip()
    fixed["level"] = str(rec.get("level") or "District").strip()
    fixed["district"] = str(rec.get("district") or "").strip()
    fixed["dob"] = str(rec.get("dob") or "").strip()
    fixed["achievements"] = str(rec.get("achievements") or "").strip()

    # Entitlements (CRITICAL: make sure always exists as list)
    ent = rec.get("entitlements_claimed", None)
    if ent is None:
        # migrate old formats:
        ent = rec.get("entitlements", None)
    if ent is None:
        ent = []
    if isinstance(ent, str):
        # allow comma-separated input
        ent = [x.strip() for x in ent.split(",") if x.strip()]
    if not isinstance(ent, list):
        ent = []
    fixed["entitlements_claimed"] = sorted(list({str(x).strip() for x in ent if str(x).strip()}))

    # Evidence fields (optional but helps analysis)
    fixed["docs_submitted"] = rec.get("docs_submitted", {})
    if not isinstance(fixed["docs_submitted"], dict):
        fixed["docs_submitted"] = {}

    fixed["reported_barriers"] = rec.get("reported_barriers", [])
    if isinstance(fixed["reported_barriers"], str):
        fixed["reported_barriers"] = [x.strip() for x in fixed["reported_barriers"].split(",") if x.strip()]
    if not isinstance(fixed["reported_barriers"], list):
        fixed["reported_barriers"] = []

    # Metadata
    fixed["recorded_at"] = rec.get("recorded_at") or rec.get("created_at") or now_iso()
    return fixed


def load_athletes() -> List[Dict[str, Any]]:
    data = safe_load_json(ATHLETES_FILE, [])
    if not isinstance(data, list):
        data = []
    normalized = [normalize_athlete_record(x if isinstance(x, dict) else {}) for x in data]
    # Save back normalized to permanently fix old messy rows
    safe_save_json(ATHLETES_FILE, normalized)
    return normalized


def add_athlete(rec: Dict[str, Any]) -> None:
    athletes = load_athletes()
    athletes.append(normalize_athlete_record(rec))
    safe_save_json(ATHLETES_FILE, athletes)


# =========================
# FILE TEXT EXTRACTION
# =========================
def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        return file_bytes.decode("latin-1", errors="ignore")


def extract_text_from_pdf(file_bytes: bytes) -> str:
    # Prefer pdfplumber (better text)
    if pdfplumber is not None:
        try:
            text_parts = []
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    if t.strip():
                        text_parts.append(t)
            return "\n\n".join(text_parts).strip()
        except Exception:
            pass

    # Fallback PyPDF2
    if PyPDF2 is not None:
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            parts = []
            for p in reader.pages:
                parts.append(p.extract_text() or "")
            return "\n\n".join(parts).strip()
        except Exception:
            return ""

    return ""


def extract_text_from_pptx(file_bytes: bytes) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(parts).strip()
    except Exception:
        return ""


def extract_text_from_xlsx(file_bytes: bytes) -> Tuple[str, Dict[str, pd.DataFrame]]:
    # Returns (combined_text, sheets_as_df)
    xls = pd.ExcelFile(io.BytesIO(file_bytes))
    sheets = {}
    text_parts = []
    for name in xls.sheet_names:
        df = xls.parse(name)
        sheets[name] = df
        # Create text summary from first rows
        preview = df.head(20).astype(str).fillna("")
        text_parts.append(f"[Sheet: {name}]\n" + "\n".join([" | ".join(row) for row in preview.values.tolist()]))
    return ("\n\n".join(text_parts).strip(), sheets)


def extract_text_by_type(filename: str, file_bytes: bytes):
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return {"type": "pdf", "text": extract_text_from_pdf(file_bytes), "tables": {}}
    if lower.endswith(".txt"):
        return {"type": "txt", "text": extract_text_from_txt(file_bytes), "tables": {}}
    if lower.endswith(".pptx"):
        return {"type": "pptx", "text": extract_text_from_pptx(file_bytes), "tables": {}}
    if lower.endswith(".xlsx"):
        text, sheets = extract_text_from_xlsx(file_bytes)
        return {"type": "xlsx", "text": text, "tables": sheets}
    return {"type": "unknown", "text": "", "tables": {}}


# =========================
# POLICY ANALYSIS (present + precision + enforceability)
# =========================
def clean_text(t: str) -> str:
    t = re.sub(r"\s+", " ", t or "").strip()
    return t


def score_precision(text: str, hits: int) -> int:
    """
    0-2 scale: precision means 'specificity'
    Heuristic:
      +1 if text has action verbs: shall/must/required
      +1 if text has numbers/dates/within X days
    """
    if not text:
        return 0
    t = text.lower()

    score = 0
    if re.search(r"\b(shall|must|required|mandatory|will)\b", t):
        score += 1
    if re.search(r"\b(within|days|months|years|deadline|time limit|period)\b", t) or re.search(r"\b\d+\b", t):
        score += 1
    score = min(score, 2)

    # if hits are extremely low, cap it
    if hits <= 0:
        score = 0
    return score


def score_enforceability(text: str) -> int:
    """
    0-1 scale: enforceability means "what happens if someone does NOT comply?"
    +1 if text mentions penalty/audit/inspection/disciplinary/non-compliance.
    """
    if not text:
        return 0
    t = text.lower()
    if re.search(r"\b(penalty|disciplinary|audit|inspection|non-compliance|punishment|fine)\b", t):
        return 1
    return 0


def analyze_policy_text(policy_text: str) -> pd.DataFrame:
    """
    Builds obligation matrix:
      Present: True/False
      Precision: 0-2
      Enforceability: 0-1
      Score: Present(1) + Precision + Enforceability => 0..4
    """
    raw = (policy_text or "")
    t = raw.lower()

    rows = []
    for ob in OBLIGATIONS:
        hits = 0
        for kw in ob["keywords"]:
            if kw in t:
                hits += 1

        present = hits > 0
        precision = score_precision(raw, hits) if present else 0
        enforce = score_enforceability(raw) if present else 0

        score = (1 if present else 0) + precision + enforce

        rows.append({
            "Obligation": ob["obligation"],
            "Present": present,
            "Precision (0-2)": precision,
            "Enforceability (0-1)": enforce,
            "Score (0-4)": score,
            "Why it matters": ob["why"],
            "Keyword hits": hits,
        })

    return pd.DataFrame(rows)


def governance_strength_interpretation(total_score: int, max_score: int) -> Tuple[str, str]:
    """
    Returns (label, explanation)
    """
    pct = (total_score / max_score) if max_score else 0.0
    if pct >= 0.75:
        return ("STRONG", "Policy is relatively complete and enforceable. Remaining work: digitization + measurable service levels + audit trails.")
    if pct >= 0.45:
        return ("MODERATE", "Policy exists but has gaps. Athletes can still be denied due to vague criteria, weak timelines, or missing grievance remedies.")
    return ("WEAK", "High governance failure risk. Policy text lacks clarity/enforcement; denial and delays become normal, and athletes cannot appeal effectively.")


# =========================
# ATHLETE IMPACT SIMULATION
# =========================
def compute_exclusion_risk(athletes: List[Dict[str, Any]], matrix_df: pd.DataFrame) -> pd.DataFrame:
    """
    Creates a simple risk model:
    - If Benefits & Entitlements obligation score is low and athlete claims entitlements -> risk increases
    - If Grievance is low and athlete reports barriers -> risk increases
    - If Digitization is low -> risk increases
    Output: risk score 0..100 and flags.
    """
    if matrix_df is None or matrix_df.empty:
        return pd.DataFrame()

    score_map = dict(zip(matrix_df["Obligation"], matrix_df["Score (0-4)"]))
    ben = score_map.get("Benefits & Entitlements", 0)
    gri = score_map.get("Grievance, Appeal & Remedies", 0)
    dig = score_map.get("Digitization & Transparency", 0)
    tim = score_map.get("Timelines & Service Levels", 0)

    rows = []
    for a in athletes:
        a = normalize_athlete_record(a)  # extra safety
        ent = a.get("entitlements_claimed", [])
        barriers = a.get("reported_barriers", [])

        risk = 10  # base
        flags = []

        # Entitlements mismatch
        if ent and ben <= 1:
            risk += 30
            flags.append("Entitlements claimed but policy entitlement clause is weak/unclear")

        # No grievance = no remedy
        if (barriers or ent) and gri <= 1:
            risk += 20
            flags.append("Weak grievance/appeal mechanism → denial cannot be challenged")

        # Delay risk
        if tim <= 1:
            risk += 15
            flags.append("Weak/absent timelines → delay risk")

        # Transparency risk
        if dig <= 1:
            risk += 15
            flags.append("Low digitization → record loss / non-traceability risk")

        # Barrier-based bump
        if barriers:
            risk += min(10, 2 * len(barriers))
            flags.append("Reported barriers increase exclusion likelihood")

        risk = max(0, min(100, risk))

        label = "LOW"
        if risk >= 70:
            label = "HIGH"
        elif risk >= 40:
            label = "MEDIUM"

        rows.append({
            "Athlete": a.get("athlete_name", ""),
            "Sport": a.get("sport", ""),
            "Level": a.get("level", ""),
            "District": a.get("district", ""),
            "Entitlements claimed": ", ".join(ent) if ent else "",
            "Reported barriers": ", ".join(barriers) if barriers else "",
            "Exclusion Risk (0-100)": risk,
            "Risk Level": label,
            "Why (flags)": " | ".join(flags) if flags else "No major red flags detected",
            "Recorded At": a.get("recorded_at", ""),
        })

    return pd.DataFrame(rows).sort_values(by="Exclusion Risk (0-100)", ascending=False)


# =========================
# OPTIONAL: Live sports data (SAFE)
# =========================
def safe_fetch_sports_data(query_name: str) -> Dict[str, Any]:
    """
    SAFE: Won't crash even if API fails.
    By default uses TheSportsDB v1 FREE endpoint (can fail sometimes).
    You can replace with your own API later.
    """
    query_name = (query_name or "").strip()
    if not query_name:
        return {"error": "Enter a player/team name"}

    # TheSportsDB sometimes changes/limits endpoints; keep it fully protected
    url = "https://www.thesportsdb.com/api/v1/json/1/searchplayers.php"
    try:
        r = requests.get(url, params={"p": query_name}, timeout=12)
        if r.status_code != 200:
            return {"error": f"API failed (HTTP {r.status_code}). This section is optional."}
        data = r.json()
        if not data or not data.get("player"):
            return {"error": "No player found / API returned empty."}
        # show only a few safe fields
        p = data["player"][0]
        return {
            "name": p.get("strPlayer"),
            "sport": p.get("strSport"),
            "nationality": p.get("strNationality"),
            "team": p.get("strTeam"),
            "position": p.get("strPosition"),
        }
    except Exception as e:
        return {"error": f"Live data not available right now. (Details: {str(e)[:80]}...)"}


# =========================
# STATE INIT
# =========================
if "policy_text" not in st.session_state:
    st.session_state.policy_text = ""

if "policy_sources" not in st.session_state:
    st.session_state.policy_sources = []  # list of dicts

if "matrix_df" not in st.session_state:
    # default empty matrix (filled when policy uploaded)
    st.session_state.matrix_df = pd.DataFrame()

if "athletes" not in st.session_state:
    st.session_state.athletes = load_athletes()


# =========================
# HEADER
# =========================
st.title("🏛️ " + APP_TITLE)
st.caption("From policy text → compliance gaps → athlete exclusion risk (dynamic analytics model)")

st.markdown(
    """
<div class="small-note">
<b>What this app does:</b>
<ul>
<li><b>Policy Upload:</b> Upload any policy/regulation file (PDF/XLSX/PPTX/TXT) and extract clauses.</li>
<li><b>Compliance Matrix:</b> Automatically scores obligations (Presence, Precision, Enforceability) and builds a Governance Strength Index.</li>
<li><b>Athlete Reality:</b> Structured athlete-ground data (entitlements claimed + barriers faced) stored cleanly.</li>
<li><b>Impact Simulation:</b> Converts policy gaps + athlete reality into an <b>Exclusion Risk Score</b> and interpretable reasons.</li>
</ul>
</div>
""",
    unsafe_allow_html=True
)

tabs = st.tabs([
    "Athlete Reality (Structured Data)",
    "Policy Upload (PDF/XLSX/PPTX/TXT)",
    "Compliance Matrix (Auto + Editable)",
    "Athlete Impact Simulation",
    "Failure Diagnosis (What to Reform)",
    "Optional: Live Sports Data"
])


# =========================
# TAB 1: Athlete Reality
# =========================
with tabs[0]:
    st.subheader("Athlete Ground Reality (Why this matters)")
    st.write(
        """
**Purpose:** This section captures the *real-world athlete experience* in a structured format.
If the athlete data is unstructured/scattered, governance failure cannot be measured.
So we force ONE clean data schema (no broken columns, no missing keys).

**Outcome:** Your athlete table becomes a clean dataset you can defend in reports, presentations, and evaluation.
"""
    )

    c1, c2 = st.columns([1.05, 1.45], gap="large")

    with c1:
        st.markdown("### Add Athlete Record (Structured Form)")
        athlete_name = st.text_input("Athlete Name", placeholder="e.g., Saurabh")
        sport = st.text_input("Sport", placeholder="e.g., Cricket")
        level = st.selectbox("Level", ["District", "State", "National", "International"], index=0)
        district = st.text_input("District/City", placeholder="e.g., Mumbai")
        dob = st.text_input("DOB (YYYY-MM-DD)", placeholder="e.g., 2002-11-20")
        achievements = st.text_area("Achievements / Notes", placeholder="e.g., District Captain; 2x gold medals")

        st.markdown("### What athlete claims they should get (Entitlements)")
        entitlements = st.multiselect("Entitlements Claimed", ENTITLEMENT_CHOICES, default=[])

        st.markdown("### What athlete says is blocking them (Barriers)")
        barriers = st.multiselect(
            "Barriers Reported",
            ["No ID/Registration", "Documents rejected", "Delays", "No transparency", "Favoritism", "No grievance option", "Funds not released", "Hostel denied"],
            default=[]
        )

        docs = {
            "identity_proof": st.checkbox("Identity proof submitted"),
            "age_proof": st.checkbox("Age proof submitted"),
            "performance_certificates": st.checkbox("Performance certificates submitted"),
            "income_certificate": st.checkbox("Income certificate submitted (if required)"),
            "medical_fitness": st.checkbox("Medical fitness submitted"),
        }

        if st.button("Save Athlete", type="primary", use_container_width=True):
            add_athlete({
                "athlete_name": athlete_name,
                "sport": sport,
                "level": level,
                "district": district,
                "dob": dob,
                "achievements": achievements,
                "entitlements_claimed": entitlements,
                "reported_barriers": barriers,
                "docs_submitted": docs,
                "recorded_at": now_iso(),
            })
            st.session_state.athletes = load_athletes()
            st.success("Saved successfully (structured schema enforced).")

    with c2:
        st.markdown("### Athlete Records Table (Clean Dataset)")
        st.session_state.athletes = load_athletes()
        df = pd.DataFrame(st.session_state.athletes)

        if df.empty:
            st.info("No athlete data yet. Add some records on the left.")
        else:
            view_cols = [
                "athlete_name", "sport", "level", "district",
                "dob", "achievements", "entitlements_claimed",
                "reported_barriers", "recorded_at"
            ]
            df_view = df[view_cols].copy()
            df_view["entitlements_claimed"] = df_view["entitlements_claimed"].apply(lambda x: ", ".join(x) if isinstance(x, list) else "")
            df_view["reported_barriers"] = df_view["reported_barriers"].apply(lambda x: ", ".join(x) if isinstance(x, list) else "")
            st.dataframe(df_view, use_container_width=True, height=420)

            st.download_button(
                "Download Athlete Dataset (CSV)",
                data=df_view.to_csv(index=False).encode("utf-8"),
                file_name="athletes_structured.csv",
                mime="text/csv",
                use_container_width=True
            )


# =========================
# TAB 2: Policy Upload
# =========================
with tabs[1]:
    st.subheader("Policy Upload → Text + Tables Extraction (Live Analysis)")
    st.write(
        """
Upload any **PDF / XLSX / PPTX / TXT** policy or compliance document.
The app extracts text and then automatically detects governance obligations.

✅ This is NOT restricted to any fixed policy list.  
✅ Whatever you upload becomes the live source for compliance scoring.
"""
    )

    uploaded = st.file_uploader(
        "Upload policy/regulation documents (multiple allowed)",
        type=["pdf", "xlsx", "pptx", "txt"],
        accept_multiple_files=True
    )

    extracted_texts = []
    extracted_tables = {}

    if uploaded:
        for f in uploaded:
            b = f.read()
            parsed = extract_text_by_type(f.name, b)
            st.session_state.policy_sources.append({
                "filename": f.name,
                "type": parsed["type"],
                "uploaded_at": now_iso(),
                "text_len": len(parsed["text"] or ""),
            })
            if parsed["text"]:
                extracted_texts.append(f"\n\n===== FILE: {f.name} =====\n{parsed['text']}\n")
            if parsed["tables"]:
                extracted_tables[f.name] = parsed["tables"]

        combined = "\n".join(extracted_texts).strip()
        st.session_state.policy_text = combined

        st.success(f"Extracted text from {len(uploaded)} file(s).")

    colA, colB = st.columns([1.1, 1.0], gap="large")
    with colA:
        st.markdown("### Extracted Policy Text (Preview)")
        preview = st.session_state.policy_text[:12000] if st.session_state.policy_text else ""
        st.text_area("Text", value=preview, height=420)

        st.download_button(
            "Download Extracted Policy Text (TXT)",
            data=(st.session_state.policy_text or "").encode("utf-8"),
            file_name="extracted_policy_text.txt",
            mime="text/plain",
            use_container_width=True
        )

    with colB:
        st.markdown("### Extracted Tables (from XLSX)")
        if not extracted_tables:
            st.info("Upload an XLSX to see sheet tables here.")
        else:
            for fname, sheets in extracted_tables.items():
                st.write(f"**{fname}**")
                for sheet_name, df_sheet in sheets.items():
                    st.caption(f"Sheet: {sheet_name}")
                    st.dataframe(df_sheet.head(30), use_container_width=True, height=240)

    st.markdown("---")
    st.markdown("### Auto-Analyze Policy (Build Compliance Matrix)")
    if st.button("Analyze Policy Now", type="primary"):
        if not st.session_state.policy_text.strip():
            st.error("No policy text available. Upload a file first.")
        else:
            st.session_state.matrix_df = analyze_policy_text(st.session_state.policy_text)
            st.success("Policy analyzed. Go to 'Compliance Matrix' tab.")


# =========================
# TAB 3: Compliance Matrix
# =========================
with tabs[2]:
    st.subheader("Compliance Obligation Matrix (Meaning + Logic)")
    st.write(
        """
This table converts policy text into a **governance measurement model**.

**How scores work (per obligation):**
- **Present (0/1):** Is the obligation mentioned at all?
- **Precision (0–2):** Are rules specific (shall/must + numbers/timelines)?
- **Enforceability (0–1):** Are penalties/audit/non-compliance consequences defined?
- **Score (0–4):** Present + Precision + Enforceability

Then all obligation scores sum into the **Governance Strength Index**.
"""
    )

    if st.session_state.matrix_df is None or st.session_state.matrix_df.empty:
        st.info("No matrix yet. Upload a policy and click 'Analyze Policy Now'.")
    else:
        matrix = st.session_state.matrix_df.copy()

        st.markdown("### Auto-detected Matrix (Editable if you want)")
        edited = st.data_editor(
            matrix[["Obligation", "Present", "Precision (0-2)", "Enforceability (0-1)", "Score (0-4)", "Why it matters"]],
            use_container_width=True,
            hide_index=True,
            disabled=["Obligation", "Why it matters"],  # Lock key meaning text
        )

        # Recompute governance index from edited values (safe)
        try:
            total = int(pd.to_numeric(edited["Score (0-4)"], errors="coerce").fillna(0).sum())
        except Exception:
            total = int(pd.to_numeric(matrix["Score (0-4)"], errors="coerce").fillna(0).sum())

        max_score = len(OBLIGATIONS) * 4
        label, explanation = governance_strength_interpretation(total, max_score)

        st.markdown("### Governance Strength Index (Explainable)")
        mc1, mc2, mc3 = st.columns([1, 1, 2], gap="large")
        with mc1:
            st.markdown(f"<div class='metric-card'><b>Index</b><br><span style='font-size:34px'>{total}</span><br><span class='muted'>out of {max_score}</span></div>", unsafe_allow_html=True)
        with mc2:
            color = "good" if label == "STRONG" else ("warn" if label == "MODERATE" else "danger")
            st.markdown(f"<div class='metric-card'><b>Status</b><br><span class='{color}' style='font-size:28px'>{label}</span></div>", unsafe_allow_html=True)
        with mc3:
            st.markdown(f"<div class='metric-card'><b>Interpretation</b><br>{explanation}</div>", unsafe_allow_html=True)

        st.download_button(
            "Download Compliance Matrix (CSV)",
            data=matrix.to_csv(index=False).encode("utf-8"),
            file_name="compliance_matrix.csv",
            mime="text/csv",
            use_container_width=True
        )


# =========================
# TAB 4: Athlete Impact Simulation
# =========================
with tabs[3]:
    st.subheader("Athlete–Policy Stress Test (Who gets excluded and why)")
    st.write(
        """
**Purpose:** Convert *policy gaps* + *athlete realities* into a **measurable exclusion risk**.

**What Exclusion Risk (0–100) means:**
- **0–39 (LOW):** Governance text + systems likely support athletes.
- **40–69 (MEDIUM):** Some gaps; athletes may be delayed/denied.
- **70–100 (HIGH):** Governance failure likely; benefits are claimed but not protected by enforceable policy.

This is where you demonstrate **real-world governance failure** analytically.
"""
    )

    if st.session_state.matrix_df is None or st.session_state.matrix_df.empty:
        st.warning("No policy matrix found. Upload policy → analyze it first.")
    else:
        athletes = load_athletes()
        if not athletes:
            st.info("No athletes added yet. Add athlete records in Athlete Reality tab.")
        else:
            impact_df = compute_exclusion_risk(athletes, st.session_state.matrix_df)

            st.markdown("### Exclusion Risk Table (Explainable)")
            st.dataframe(impact_df, use_container_width=True, height=420)

            st.download_button(
                "Download Exclusion Risk Report (CSV)",
                data=impact_df.to_csv(index=False).encode("utf-8"),
                file_name="exclusion_risk_report.csv",
                mime="text/csv",
                use_container_width=True
            )


# =========================
# TAB 5: Failure Diagnosis
# =========================
with tabs[4]:
    st.subheader("Failure Diagnosis & Reform Plan (Actionable)")
    st.write(
        """
This section converts your analytics into **reform language**:
- What obligations are missing?
- What does it cause on ground?
- What should government/university/sports bodies change?

This is the “presentable” layer for your model.
"""
    )

    if st.session_state.matrix_df is None or st.session_state.matrix_df.empty:
        st.info("Upload policy → analyze it to generate diagnosis.")
    else:
        matrix = st.session_state.matrix_df.copy()
        weak = matrix.sort_values("Score (0-4)").head(5)

        st.markdown("### Top Governance Gaps (Lowest scoring obligations)")
        st.dataframe(weak[["Obligation", "Score (0-4)", "Why it matters"]], use_container_width=True, hide_index=True)

        st.markdown("### Reform Checklist (What to add to policy)")
        checklist = []
        for _, r in weak.iterrows():
            obligation = r["Obligation"]
            score = int(r["Score (0-4)"])
            if score <= 1:
                checklist.append({
                    "Reform item": f"Strengthen '{obligation}' clauses",
                    "Minimum requirement": "Define process + documents + timelines + responsible authority",
                    "Add enforceability": "Include audit/penalty/non-compliance consequences",
                    "Why": r["Why it matters"],
                })

        if checklist:
            chk_df = pd.DataFrame(checklist)
            st.dataframe(chk_df, use_container_width=True, height=320)

            st.download_button(
                "Download Reform Checklist (CSV)",
                data=chk_df.to_csv(index=False).encode("utf-8"),
                file_name="reform_checklist.csv",
                mime="text/csv",
                use_container_width=True
            )
        else:
            st.success("No critical gaps detected (policy appears strong in this model).")

        st.markdown("---")
        st.markdown("### Quick Narrative (Use in your report/presentation)")
        total = int(pd.to_numeric(matrix["Score (0-4)"], errors="coerce").fillna(0).sum())
        max_score = len(OBLIGATIONS) * 4
        label, explanation = governance_strength_interpretation(total, max_score)

        narrative = f"""
**Governance Strength Index:** {total}/{max_score} → **{label}**  
**Interpretation:** {explanation}

**Key failure drivers (lowest scoring obligations):**
{", ".join(weak["Obligation"].tolist())}

**Expected ground impact:** weak entitlement clarity + weak grievance + weak timelines leads to denial, delay, and non-traceability—creating systematic athlete exclusion.
"""
        st.text_area("Narrative", value=narrative.strip(), height=200)


# =========================
# TAB 6: Optional Live Sports Data
# =========================
with tabs[5]:
    st.subheader("Optional: Live Sports Data (Will NOT break your model)")
    st.write(
        """
This is optional. Your core model is governance failure + compliance analytics.
If live data API fails, it should never crash the app—so we keep it isolated and safe.
"""
    )

    name = st.text_input("Search player name (optional)", placeholder="e.g., Rohit Sharma")
    if st.button("Fetch Live Data"):
        data = safe_fetch_sports_data(name)
        st.json(data)

    st.caption("If you want guaranteed live data, we can plug in a paid/official API later (but the governance model does not depend on it).")
